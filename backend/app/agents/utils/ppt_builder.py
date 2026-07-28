import os
import logging
from typing import List, Optional, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import subprocess
import tempfile

logger = logging.getLogger(__name__)

class PowerPointBuilder:
    '''Builder for creating PowerPoint presentations'''
    
    def __init__(self, output_dir: str = "./generated/ppt"):
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)
        self.prs = None
    
    def create_presentation(self, title: str, theme: str = "modern") -> Presentation:
        '''Create a new presentation'''
        self.prs = Presentation()
        
        # Set slide dimensions (16:9)
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        return self.prs
    
    def add_title_slide(self, title: str, subtitle: str = "") -> None:
        '''Add a title slide'''
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        if subtitle and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
    
    def add_content_slide(
        self,
        title: str,
        content: str = "",
        bullet_points: List[str] = None,
        speaker_notes: str = "",
        images: List[str] = None,
        diagrams: List[str] = None,
        videos: List[str] = None
    ) -> None:
        '''Add a content slide with various elements'''
        
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                content_placeholder = shape
                break
        
        if content_placeholder:
            if bullet_points and len(bullet_points) > 0:
                tf = content_placeholder.text_frame
                tf.text = ""
                for point in bullet_points:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
                    p.font.size = Pt(24)
            elif content:
                tf = content_placeholder.text_frame
                tf.text = content
                tf.paragraphs[0].font.size = Pt(20)
        
        if speaker_notes and hasattr(slide, 'notes_slide'):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes
        
        if images:
            for img_path in images[:1]:
                if os.path.exists(img_path):
                    try:
                        left = Inches(8)
                        top = Inches(2)
                        width = Inches(4)
                        slide.shapes.add_picture(img_path, left, top, width=width)
                    except Exception as e:
                        logger.warning(f"Failed to add image {img_path}: {e}")
        
        if diagrams:
            for diagram_path in diagrams[:1]:
                if os.path.exists(diagram_path):
                    try:
                        left = Inches(1)
                        top = Inches(2)
                        width = Inches(6)
                        slide.shapes.add_picture(diagram_path, left, top, width=width)
                    except Exception as e:
                        logger.warning(f"Failed to add diagram {diagram_path}: {e}")
    
    def add_diagram_slide(
        self,
        title: str,
        diagram_path: str,
        description: str = ""
    ) -> None:
        '''Add a slide with a diagram'''
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        if os.path.exists(diagram_path):
            try:
                left = Inches(1)
                top = Inches(1.5)
                width = Inches(11)
                slide.shapes.add_picture(diagram_path, left, top, width=width)
            except Exception as e:
                logger.warning(f"Failed to add diagram: {e}")
        
        if description:
            try:
                text_box = slide.shapes.add_textbox(
                    Inches(1), Inches(6.5), Inches(11), Inches(0.8)
                )
                tf = text_box.text_frame
                tf.text = description
                tf.paragraphs[0].font.size = Pt(16)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            except Exception as e:
                logger.warning(f"Failed to add description: {e}")
    
    def add_video_slide(
        self,
        title: str,
        video_path: str,
        description: str = ""
    ) -> None:
        '''Add a slide with a video link/embedding'''
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        try:
            left = Inches(2)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5)
            
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(50, 50, 50)
            shape.line.color.rgb = RGBColor(100, 100, 100)
            
            text_box = slide.shapes.add_textbox(
                left + Inches(1), top + Inches(2), width - Inches(2), Inches(1)
            )
            tf = text_box.text_frame
            tf.text = "VIDEO"
            tf.paragraphs[0].font.size = Pt(36)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
            
            if video_path:
                text_box2 = slide.shapes.add_textbox(
                    left + Inches(1), top + Inches(3.5), width - Inches(2), Inches(0.5)
                )
                tf2 = text_box2.text_frame
                tf2.text = f"Video: {os.path.basename(video_path)}"
                tf2.paragraphs[0].font.size = Pt(14)
                tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
                tf2.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
            
        except Exception as e:
            logger.warning(f"Failed to add video placeholder: {e}")
        
        if description:
            try:
                text_box3 = slide.shapes.add_textbox(
                    Inches(1), Inches(6.5), Inches(11), Inches(0.8)
                )
                tf3 = text_box3.text_frame
                tf3.text = description
                tf3.paragraphs[0].font.size = Pt(16)
                tf3.paragraphs[0].alignment = PP_ALIGN.CENTER
            except Exception as e:
                logger.warning(f"Failed to add description: {e}")
    
    def save(self, filename: str) -> str:
        '''Save the presentation'''
        if not filename.endswith('.pptx'):
            filename += '.pptx'
        
        output_path = os.path.join(self.output_dir, filename)
        self.prs.save(output_path)
        logger.info(f"Presentation saved to: {output_path}")
        return output_path
