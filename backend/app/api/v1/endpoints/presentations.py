import asyncio
import json
import logging
import os
import re
import subprocess
import tempfile
import uuid
from pathlib import Path
from typing import Dict, List, Optional

from fastapi import APIRouter, BackgroundTasks, HTTPException
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt

from app.services.animation_generator import AnimationGenerator
from app.services.diagram_generator import DiagramGenerator
from app.services.llm import LLMService
from app.services.settings_service import settings_service
from app.services.timeline_service import timeline_service

router = APIRouter()
logger = logging.getLogger(__name__)

class GenerateRequest(BaseModel):
    prompt: str
    slide_count: Optional[int] = 10
    audience: Optional[str] = "general"
    tone: Optional[str] = "educational"

generated_content = {}


def _existing_diagrams(job_id: str) -> List[Dict]:
    diagrams = generated_content.get(job_id, {}).get("diagrams", [])
    return [
        diagram
        for diagram in diagrams
        if diagram.get("image_path") and os.path.exists(diagram["image_path"])
    ]


def _existing_animation(job_id: str) -> Optional[Dict]:
    animation = generated_content.get(job_id, {}).get("animation")
    if animation and animation.get("video_path") and os.path.exists(animation["video_path"]):
        return animation
    return None


def _serialize_diagrams(job_id: str) -> List[Dict]:
    serialized = []
    for index, diagram in enumerate(_existing_diagrams(job_id)):
        serialized.append(
            {
                "title": diagram.get("title", "Diagram"),
                "description": diagram.get("description", ""),
                "type": diagram.get("type", "diagram"),
                "slide_number": diagram.get("slide_number"),
                "preview_url": f"/api/v1/presentations/{job_id}/assets/diagrams/{index}",
            }
        )
    return serialized


def _serialize_animation(job_id: str) -> Dict:
    animation = _existing_animation(job_id)
    if not animation:
        return {}

    return {
        "title": animation.get("title", "Animation"),
        "description": animation.get("description", ""),
        "preview_url": f"/api/v1/presentations/{job_id}/assets/animation",
        "download_url": f"/api/v1/presentations/{job_id}/download/animation",
    }


def _diagram_map(diagrams: List[Dict]) -> Dict[int, Dict]:
    mapped = {}
    for diagram in diagrams:
        slide_number = diagram.get("slide_number")
        if slide_number and slide_number not in mapped:
            mapped[slide_number] = diagram
    return mapped


def _ensure_animation_poster(job_id: str, video_path: str) -> Optional[str]:
    if not video_path or not os.path.exists(video_path):
        return None

    poster_dir = Path(video_path).parent
    poster_path = poster_dir / "poster.png"
    if poster_path.exists():
        return str(poster_path)

    try:
        subprocess.run(
            [
                "ffmpeg",
                "-y",
                "-i",
                video_path,
                "-ss",
                "00:00:00.5",
                "-vframes",
                "1",
                str(poster_path),
            ],
            capture_output=True,
            text=True,
            timeout=30,
            check=False,
        )
    except Exception as exc:
        logger.warning("Could not create animation poster for %s: %s", job_id, exc)

    if poster_path.exists():
        return str(poster_path)
    return None

@router.post("/generate")
async def generate_presentation(request: GenerateRequest, background_tasks: BackgroundTasks):
    try:
        job_id = str(uuid.uuid4())
        logger.info(f"Generation started for job: {job_id}")
        logger.info(f"Prompt: {request.prompt[:100]}...")
        
        timeline_service.create_job(job_id, request.prompt)
        background_tasks.add_task(
            simulate_generation, 
            job_id, 
            request.prompt, 
            request.slide_count
        )
        
        return {
            "job_id": job_id,
            "status": "started",
            "message": f"Presentation generation started for: {request.prompt[:50]}..."
        }
    except Exception as e:
        logger.error(f"Generation error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

async def simulate_generation(job_id: str, prompt: str, slide_count: int = 10):
    try:
        provider = settings_service.get_setting("llm_provider", "groq")
        llm = LLMService(provider=provider)
        diagram_gen = DiagramGenerator(llm)
        animation_gen = AnimationGenerator(llm)
        generated_content[job_id] = {"slides": [], "diagrams": [], "animation": None, "title": prompt}
        
        # Step 1: Planning
        timeline_service.update_step(job_id, "planning", "in_progress", "Analyzing prompt...")
        await asyncio.sleep(0.5)
        timeline_service.update_step(job_id, "planning", "completed", "Planning complete")
        
        # Step 2: Generate content
        timeline_service.update_step(job_id, "research", "in_progress", "Generating content...")
        
        system_prompt = f"""You are a presentation expert. Create a detailed presentation outline with {slide_count} slides about: {prompt}

        Return ONLY valid JSON with this structure:
        {{
            "title": "Presentation Title",
            "slides": [
                {{
                    "number": 1,
                    "title": "Slide Title",
                    "content": "Main content for the slide",
                    "bullet_points": ["Key point 1", "Key point 2", "Key point 3"],
                    "speaker_notes": "Speaker notes for this slide"
                }}
            ],
            "key_concepts": ["Concept 1", "Concept 2", "Concept 3"]
        }}
        
        Make the content educational and informative.
        """
        
        response = await llm.generate(
            prompt=f"Create a presentation about: {prompt}",
            system_prompt=system_prompt,
            temperature=0.4,
            max_tokens=4096
        )
        
        import re
        json_match = re.search(r'\{.*\}', response.text, re.DOTALL)
        if json_match:
            content_data = json.loads(json_match.group())
        else:
            content_data = {"title": prompt, "slides": [], "key_concepts": []}
        
        timeline_service.update_step(job_id, "research", "completed", "Research complete")
        timeline_service.update_step(job_id, "slides", "completed", "Slides generated")
        generated_content[job_id]["slides"] = content_data.get("slides", [])
        generated_content[job_id]["title"] = content_data.get("title", prompt)
        
        # Step 3: Generate diagrams
        if settings_service.get_setting("enable_diagrams", True):
            timeline_service.update_step(job_id, "diagrams", "in_progress", "Creating diagrams...")
            diagrams = await diagram_gen.generate_diagrams(job_id, prompt, content_data.get("slides", []))
            timeline_service.update_step(
                job_id,
                "diagrams",
                "completed",
                f"{len(diagrams)} useful diagrams created" if diagrams else "No useful diagrams required",
            )
        else:
            diagrams = []
            timeline_service.update_step(job_id, "diagrams", "completed", "Diagrams disabled")

        generated_content[job_id]["diagrams"] = diagrams
        
        # Step 4: Generate animation
        animation = None
        if settings_service.get_setting("enable_animations", True):
            timeline_service.update_step(job_id, "animation", "in_progress", "Evaluating animation suitability...")
            animation = await animation_gen.generate_animation(job_id, prompt, content_data.get("slides", []))
            generated_content[job_id]["animation"] = animation
            timeline_service.update_step(
                job_id,
                "animation",
                "completed",
                "Animation created" if animation else "No useful animation required",
            )
        else:
            timeline_service.update_step(job_id, "animation", "completed", "Animations disabled")
        
        # Step 5: Build PowerPoint
        timeline_service.update_step(job_id, "ppt", "in_progress", "Building PowerPoint...")
        await generate_pptx(job_id, content_data, prompt, diagrams, animation)
        timeline_service.update_step(job_id, "ppt", "completed", "PowerPoint created")
        
        # Step 6: Generate PDF
        timeline_service.update_step(job_id, "pdf", "in_progress", "Exporting to PDF...")
        await generate_pdf(job_id, content_data, prompt, diagrams, animation)
        timeline_service.update_step(job_id, "pdf", "completed", "PDF exported")
        
        # Step 7: Generate Speaker Notes
        await generate_notes(job_id, content_data, prompt)
        
        # Step 8: Complete
        timeline_service.update_step(job_id, "complete", "completed", "Presentation ready!")

    except Exception as e:
        logger.error(f"Generation error: {e}")
        timeline_service.update_step(job_id, "complete", "failed", f"Error: {str(e)}")

async def generate_pptx(job_id: str, content_data: dict, prompt: str, diagrams: list, animation: Optional[Dict]):
    try:
        prs = Presentation()

        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = content_data.get("title", f"Presentation on {prompt}")
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = f"Generated by Vynetra AI\n{prompt}"

        slides_data = content_data.get("slides", [])
        if not slides_data:
            slides_data = [
                {
                    "number": 1,
                    "title": f"Introduction to {prompt}",
                    "content": f"Welcome to this presentation about {prompt}.",
                    "bullet_points": ["Key concept 1", "Key concept 2", "Key concept 3"],
                },
            ]
        diagram_map = _diagram_map(diagrams)

        for slide_data in slides_data[:8]:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.6))
            title_frame = title_box.text_frame
            title_frame.text = slide_data.get("title", f"Slide {slide_data.get('number', 0)}")
            title_frame.paragraphs[0].font.size = Pt(26)
            title_frame.paragraphs[0].font.bold = True

            diagram = diagram_map.get(slide_data.get("number"))
            content_width = Inches(5.4 if diagram else 11.6)
            content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), content_width, Inches(5.6))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            intro = text_frame.paragraphs[0]
            intro.text = slide_data.get("content", "")
            intro.font.size = Pt(18)

            for bullet in slide_data.get("bullet_points", []):
                paragraph = text_frame.add_paragraph()
                paragraph.text = bullet
                paragraph.level = 0
                paragraph.font.size = Pt(16)

            if diagram and diagram.get("image_path") and os.path.exists(diagram["image_path"]):
                slide.shapes.add_picture(diagram["image_path"], Inches(6.5), Inches(1.6), width=Inches(5.6))

        if animation and animation.get("video_path") and os.path.exists(animation["video_path"]):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.6))
            title_frame = title_box.text_frame
            title_frame.text = animation.get("title", "Animation")
            title_frame.paragraphs[0].font.size = Pt(26)
            title_frame.paragraphs[0].font.bold = True

            description = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(11.8), Inches(0.8))
            desc_frame = description.text_frame
            desc_frame.text = animation.get("description", "Visual explanation generated by Vynetra.")
            desc_frame.paragraphs[0].font.size = Pt(16)

            poster_path = _ensure_animation_poster(job_id, animation["video_path"])
            if poster_path and hasattr(slide.shapes, "add_movie"):
                slide.shapes.add_movie(
                    animation["video_path"],
                    Inches(1.0),
                    Inches(1.9),
                    Inches(8.8),
                    Inches(4.8),
                    poster_frame_image=poster_path,
                    mime_type="video/mp4",
                )
            elif poster_path:
                slide.shapes.add_picture(poster_path, Inches(1.0), Inches(1.9), width=Inches(8.8))
        
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        prs.save(temp_file.name)
        temp_file.close()
        
        generated_content[job_id]["pptx_path"] = temp_file.name
        
        logger.info(f"PPTX generated for job: {job_id}")
        
    except Exception as e:
        logger.error(f"Error generating PPTX: {e}")

async def generate_pdf(job_id: str, content_data: dict, prompt: str, diagrams: list, animation: Optional[Dict]):
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.enums import TA_CENTER, TA_LEFT
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.platypus import Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer

        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
        doc = SimpleDocTemplate(temp_file.name, pagesize=letter)
        styles = getSampleStyleSheet()

        title_style = ParagraphStyle("CustomTitle", parent=styles["Title"], fontSize=24, alignment=TA_CENTER, spaceAfter=20)
        heading_style = ParagraphStyle("CustomHeading", parent=styles["Heading1"], fontSize=16, spaceAfter=10)
        body_style = ParagraphStyle("CustomBody", parent=styles["Normal"], fontSize=11, alignment=TA_LEFT, spaceAfter=8)
        bullet_style = ParagraphStyle("Bullet", parent=styles["Normal"], fontSize=11, leftIndent=20, spaceAfter=4)

        story = []
        story.append(Paragraph(content_data.get("title", f"Presentation on {prompt}"), title_style))
        story.append(Spacer(1, 12))
        story.append(Paragraph(f"Topic: {prompt}", styles["Normal"]))
        story.append(Spacer(1, 12))
        story.append(Paragraph("Generated by Vynetra AI", styles["Normal"]))
        story.append(PageBreak())

        slides_data = content_data.get("slides", [])
        diagram_map = _diagram_map(diagrams)
        for slide_data in slides_data[:8]:
            story.append(Paragraph(f"Slide {slide_data.get('number', 0)}: {slide_data.get('title', 'Slide')}", heading_style))
            story.append(Paragraph(slide_data.get("content", ""), body_style))
            for bp in slide_data.get("bullet_points", []):
                story.append(Paragraph(f"• {bp}", bullet_style))

            diagram = diagram_map.get(slide_data.get("number"))
            if diagram and diagram.get("image_path") and os.path.exists(diagram["image_path"]):
                story.append(Spacer(1, 10))
                story.append(Image(diagram["image_path"], width=460, height=250))

            story.append(Spacer(1, 12))
            story.append(PageBreak())

        if animation and animation.get("video_path") and os.path.exists(animation["video_path"]):
            poster_path = _ensure_animation_poster(job_id, animation["video_path"])
            if poster_path and os.path.exists(poster_path):
                story.append(Paragraph(animation.get("title", "Animation"), heading_style))
                story.append(Paragraph(animation.get("description", ""), body_style))
                story.append(Image(poster_path, width=460, height=260))
                story.append(Spacer(1, 12))

        doc.build(story)
        temp_file.close()
        
        generated_content[job_id]["pdf_path"] = temp_file.name
        
        logger.info(f"PDF generated for job: {job_id}")
        
    except Exception as e:
        logger.error(f"Error generating PDF: {e}")

async def generate_notes(job_id: str, content_data: dict, prompt: str):
    try:
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".txt")
        content = []
        
        content.append("=" * 60)
        content.append(f"SPEAKER NOTES: {content_data.get('title', prompt)}")
        content.append("=" * 60)
        content.append("")
        
        slides_data = content_data.get("slides", [])
        for slide_data in slides_data[:10]:
            title = slide_data.get("title", f"Slide {slide_data.get('number', 0)}")
            notes = slide_data.get("speaker_notes", slide_data.get("content", ""))
            
            content.append("")
            content.append(f"SLIDE {slide_data.get('number', 0)}: {title}")
            content.append("-" * 40)
            content.append(notes)
            content.append("")
        
        content.append("")
        content.append("=" * 60)
        content.append(f"Generated by Vynetra AI")
        content.append(f"Job ID: {job_id}")
        content.append("=" * 60)
        
        full_content = "\n".join(content)
        temp_file.write(full_content.encode('utf-8'))
        temp_file.close()
        
        generated_content[job_id]["notes_path"] = temp_file.name
        
        logger.info(f"Notes generated for job: {job_id}")
        
    except Exception as e:
        logger.error(f"Error generating notes: {e}")

@router.get("/{job_id}/content")
async def get_content(job_id: str):
    if job_id in generated_content:
        return {
            "job_id": job_id,
            "title": generated_content[job_id].get("title", ""),
            "slides": generated_content[job_id].get("slides", []),
            "diagrams": _serialize_diagrams(job_id),
            "animation": _serialize_animation(job_id),
        }
    return {
        "job_id": job_id,
        "error": "Content not found"
    }

@router.get("/{job_id}/download/{file_type}")
async def download_file(job_id: str, file_type: str):
    try:
        if job_id in generated_content:
            if file_type == "pptx" and "pptx_path" in generated_content[job_id]:
                return FileResponse(
                    generated_content[job_id]["pptx_path"],
                    media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    filename=f"presentation_{job_id}.pptx"
                )
            elif file_type == "pdf" and "pdf_path" in generated_content[job_id]:
                return FileResponse(
                    generated_content[job_id]["pdf_path"],
                    media_type="application/pdf",
                    filename=f"presentation_{job_id}.pdf"
                )
            elif file_type == "notes" and "notes_path" in generated_content[job_id]:
                return FileResponse(
                    generated_content[job_id]["notes_path"],
                    media_type="text/plain",
                    filename=f"speaker_notes_{job_id}.txt"
                )
            elif file_type == "animation":
                anim = _existing_animation(job_id)
                if anim:
                    return FileResponse(
                        anim["video_path"],
                        media_type="video/mp4",
                        filename=f"animation_{job_id}.mp4",
                    )

        return JSONResponse(
            status_code=404,
            content={"error": "File not found"}
        )
            
    except Exception as e:
        logger.error(f"Download error: {e}")
        return JSONResponse(
            status_code=500,
            content={"error": str(e)}
        )


@router.get("/{job_id}/assets/animation")
async def preview_animation(job_id: str):
    animation = _existing_animation(job_id)
    if not animation:
        raise HTTPException(status_code=404, detail="Animation not found")

    return FileResponse(animation["video_path"], media_type="video/mp4")


@router.get("/{job_id}/assets/diagrams/{diagram_index}")
async def preview_diagram(job_id: str, diagram_index: int):
    diagrams = _existing_diagrams(job_id)
    if diagram_index < 0 or diagram_index >= len(diagrams):
        raise HTTPException(status_code=404, detail="Diagram not found")

    diagram = diagrams[diagram_index]
    return FileResponse(diagram["image_path"], media_type="image/png")

@router.get("/{job_id}/status")
async def get_status(job_id: str):
    timeline = timeline_service.get_timeline(job_id)
    if "error" in timeline:
        return {
            "job_id": job_id,
            "status": "not_found",
            "progress": 0,
            "message": "Job not found"
        }
    
    return {
        "job_id": job_id,
        "status": timeline["status"],
        "progress": timeline["progress"],
        "current_step": timeline["current_step"],
        "message": f"Processing: {timeline.get('current_step', 'starting')}"
    }

@router.get("/")
async def list_presentations():
    return {
        "presentations": [],
        "total": 0
    }
